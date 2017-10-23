#!/usr/bin/env python2
# -*- coding: utf-8 -*-
"""
Created on Thu Oct 12 16:14:01 2017

@author: yxcheng
"""


#%%
######################################################
#main    


#%%
from pptx import Presentation
prs = Presentation()
#blank_slide_layout = prs.slide_layouts[6]
#slide = prs.slides.add_slide(blank_slide_layout)
#img_path = './example/Picture4.png'
#add_pic_with_title(slide, img_path, 5, 4.5, 2.5, 'Picture4.png')
#prs.save('test2.pptx')

import os
from utils import *
img_paths = []
for file in os.listdir("./example"):
    if file.endswith(".png"):
        print(os.path.join("./example", file))
        img_paths.append(os.path.join("./example", file))



img_name = ['0','1','2','3','4','5','6','7']
img_nums = 8
prs = add_new_slide_pics(prs, img_paths, img_nums, img_name, 'Title')

add_a_table(prs, 4, 3, 'test the table')
prs.save('test1.pptx')

#import re
file_name = file_lists[0]
file_name = clean_file_name(file_name)

file_name2 = 'UWEVEREST_981__UWEVEREST-981-H2212406_19750320_Male_Angio (3mmx3mm)_20170104135911_OS_20170104144338_Angiography_Avascular.bmp'
file_name2 = clean_file_name(file_name2)

file_name3 = './example2/Angio6x6/UWANG5000_977__UWANG5000-977-H3848027_19741016_Female_Angiography 6x6 mm_20161123162131_OS_20161123163020_Structure_Superficial.bmp'
file_name3 = clean_file_name(file_name3)

try:
    found = re.search("(?<=Angiography).*?(?=mm)", file_name)
except AttributeError:
    print(' pattern not found')
#    return -1
found.group(0) 

#########
from pptx import Presentation
import re
import os
from utils import *

prs = Presentation()

file_lists = search_images_in_folder('./example2/UWANG5000-977/')
df = analysis_filelists(file_lists)

prs = add_a_table(prs, 12, 2, 'Summary', df)

#scan_time = '20161123161757'
#prs = build_the_slide(prs, df, 'Retina', scan_time)
prs = create_slides(df, prs)

prs.save('test2.pptx')
print('Slides saved!')



                
######################################################
#%% test part
from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!" 
subtitle.text = "python-pptx was here!"

prs.save('test.pptx')


#%%
from pptx import Presentation

prs = Presentation()
bullet_slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Adding a Bullet Slide'

tf = body_shape.text_frame
tf.text = 'Find the bullet slide layout'

p = tf.add_paragraph()
p.text = 'Use _TextFrame.text for first bullet'
p.level = 1

p = tf.add_paragraph()
p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
p. level = 2

prs.save('test.pptx')
#%%
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = width = height = Inches(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame

tf.text = "This is text inside a textbox"

p = tf.add_paragraph()
p.text = "This is a second paragraph that's bold"
p.font.bold = True

p = tf.add_paragraph()
p.text = "This is a third paragraph that's big"
p.font.size = Pt(40)

prs.save('test.pptx')
#%%
from pptx import Presentation
from pptx.util import Inches



prs = Presentation() 
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

img_path = './example/Picture1.png'

left = top = Inches(1)
pic = slide.shapes.add_picture(img_path, left, top)

left = top = width = height = Inches(1)
txBox = slide.shapes.add_textbox(left, top*0.7 , width, height)
tf = txBox.text_frame

tf.text = img_path[10::]


img_path = './example/Picture2.png'
left = Inches(5)
height = Inches(2.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

top = width = height = Inches(1)
txBox = slide.shapes.add_textbox(left, top*0.7 , width, height)
tf = txBox.text_frame

tf.text = img_path[10::]
#%%
img_path = './example/Picture2.png'
left = Inches(5)
height = Inches(2.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

top = width = height = Inches(1)
txBox = slide.shapes.add_textbox(left, top*0.7 , width, height)
tf = txBox.text_frame

tf.text = img_path[10::]
    
#%% main
import argparse
if __name__ == "__main__":
    args = parse_args()
    
    