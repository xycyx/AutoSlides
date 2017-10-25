#!/usr/bin/env python2
# -*- coding: utf-8 -*-
"""
Created on Sun Oct 22 17:48:08 2017

@author: yxcheng
"""

#!/usr/bin/env python2
# -*- coding: utf-8 -*-
"""
Created on Thu Oct 12 16:14:01 2017

@author: yxcheng
"""
import os
import re
from pptx.util import Inches, Pt

#%%
def add_pic_with_title(slide, img_path, left, top, width, img_name):
    """add a picture to slide from img_path and add the title with img_name
    """
    from pptx.util import Inches
    left = Inches(left)
    top_img = Inches(top)
    width = Inches(width)
    
    # add picture
    try:
        pic = slide.shapes.add_picture(img_path, left, top_img, width=width)
    except IOError:
        print('warrning: IO error' ) 
    except AttributeError:
        print('warrning: AttributeError' )
    #add title
    top_title = Inches(top - 0.3)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top_title , width, height)
    tf = txBox.text_frame
    tf.text = img_name
    
    return slide
    
def add_new_slide_pics(prs, img_paths, img_nums, img_name, slide_title=''):
    """Decide the layout of pictures
    """
#    from pptx.util import Inches
    #img_nums = len(img_paths)

    #%% 2
    if img_nums == 2:
        layout = [1, 2]
        prs = set_layout(prs, layout, img_paths, img_nums, img_name, slide_title)
    
    #%% 3
    elif img_nums == 3:
        layout = [1, 3]
        prs = set_layout(prs, layout, img_paths, img_nums, img_name, slide_title)
   
    #%% 4 still have a issue
    elif img_nums == 4:
        layout = [2, 2]
        prs = set_layout(prs, layout, img_paths, img_nums, img_name, slide_title)
    
    #%% 6
    elif img_nums == 6:
        layout = [2, 3]
        prs = set_layout(prs, layout, img_paths, img_nums, img_name, slide_title)
    
    #%% 8
    elif img_nums == 8:
        layout = [2, 4]
        prs = set_layout(prs, layout, img_paths, img_nums, img_name, slide_title)        
               
    else:
        print('other numbers')
    return prs
    
def set_layout(prs, layout, img_paths, img_nums, img_name, slide_title=''):
    """Set the position and decide the layout of picture 
    Args: 
        layout: the size of images matriax , as [4, 2]
    Return:
        pre: the Presentation 
    """
    from pptx.util import Inches
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    txleft = txtop = txwidth = txheight = Inches(1)
    txBox = slide.shapes.add_textbox(txleft, txtop, txwidth, txheight)
    tf = txBox.text_frame
    
    tf.text = slide_title
#    p = tf.add_paragraph()
#    p.text = slide_title
#    p.font.size = Pt(22)
    
    ind = 0
    int_width = 0.3
    width = (10 - int_width)/layout[1]
    pic_width = width - int_width
    int_top = 7.5 - layout[0]*width
#    for img_path in img_paths: 
    for i in range(0, img_nums):
        img_path = img_paths[i]          
        col = ind%layout[1]
        row = ind//layout[1]
#            height = ((10 - int_width)/layout[1] - int_width
        left = col*width + int_width
        top = row*width + int_top
        
        # add the picture on its position
        add_pic_with_title(slide, img_path, left, top, pic_width, img_name[ind])
        ind = ind + 1
    return prs
#%%   
# # get the size of image without open it
#from PIL import Image
#im = Image.open('./example/1.tiff')
#im.size


#width, height = get_image_size('./example/1.jpg')
#import struct
#import imghdr
#
#def get_image_size(fname):
#    '''Determine the image type of fhandle and return its size.
#    from draco'''
#    with open(fname, 'rb') as fhandle:
#        head = fhandle.read(24)
#        if len(head) != 24:
#            return
#        if imghdr.what(fname) == 'png':
#            check = struct.unpack('>i', head[4:8])[0]
#            if check != 0x0d0a1a0a:
#                return
#            width, height = struct.unpack('>ii', head[16:24])
#        elif imghdr.what(fname) == 'gif':
#            width, height = struct.unpack('<HH', head[6:10])
#        elif imghdr.what(fname) == 'bmp':
#            _, width, height, depth = re.search(
#                b"((\d+)\sx\s"
#                b"(\d+)\sx\s"
#                b"(\d+))", str).groups()
#            width = int(width)
#            height = int(height)
#        elif imghdr.what(fname) == 'jpeg':
#            try:
#                fhandle.seek(0) # Read 0xff next
#                size = 2
#                ftype = 0
#                while not 0xc0 <= ftype <= 0xcf:
#                    fhandle.seek(size, 1)
#                    byte = fhandle.read(1)
#                    while ord(byte) == 0xff:
#                        byte = fhandle.read(1)
#                    ftype = ord(byte)
#                    size = struct.unpack('>H', fhandle.read(2))[0] - 2
#                # We are at a SOFn block
#                fhandle.seek(1, 1)  # Skip `precision' byte.
#                height, width = struct.unpack('>HH', fhandle.read(4))
#            except Exception: #IGNORE:W0703
#                return
#        else:
#            return
#        return width, height    

#######################################################################

def search_images_in_folder(directory):
    """Search and list all images (bmp, png, jpg, tif) files in the folder
    Args: 
        directory: the folder conatins all files
    Return:
        file_lists: the list of all images files
    """
    import os
    file_lists = []
#    directory = './example2'
    num = 0
    for root, dirs, files in os.walk(directory):
        for file in files:
            if file.lower().endswith(('.bmp','.png','.jpg','.jpeg','.tiff','.tif')):
#                print file
                file_lists.append(os.path.join(root, file))
    return file_lists       
 
#file_lists = search_images_in_folder('./example2')

#%%            



def clean_file_name(file_name):
    """Remove the ' ', '(', ')' of the the filename
    Args: 
        file_name: the input filename
    Returns: 
        file_name: the output of cleaned filename
    """
    
    while ' ' in file_name:
        file_name = file_name.replace(' ', '')
    while '(' in file_name:
        file_name = file_name.replace('(', '')
    while ')' in file_name:
        file_name = file_name.replace(')', '')           
    return file_name

#%%
# regular expression to parse the file name
def find_system_type(file_name):
    """ get system type by the rule 'UWANG5000' or 'UWEVEREST'
    """
    if file_name.upper().find('UWANG5000')>=0:
        system_type = 'SD'
    elif file_name.upper().find('UWEVEREST')>=0:
        system_type = 'SS'
    else:
        print ('system type not found')
        system_type = -1
    return system_type

def find_first_name(file_name):
    """ get the first name by the rule '_XXXX__'
    """
    # the the last occurrence of '/', as the strating position of firstname
    if file_name.rfind('/')>=0:
        name_position = file_name.rfind('/') + 1
    else:
        name_position = 0
    
    #get the first name from the strat position    
    try:
        found = re.search(".+[_][_]", file_name[name_position:])
        first_name_str = found.group(0)
    except AttributeError:
        print('first name number not found')
        return -1
    first_name = first_name_str[:-2]
    return first_name

def find_case_number(file_name):
    """ get the case number by rule '-number-'
    """
    case_str = find_first_name(file_name)
    try:
        found = re.search("[-|_]\d+", case_str)
        case_str_number = found.group(0)
    except AttributeError:
        print('case number not found')
        return -1
    case_number = int(case_str_number[1:])
    return case_number

def find_gender(file_name):
    """ get system type by '_male_' or '_female_'
    """
    if file_name.lower().find('_female_')>=0:
        gender = 'female'
    elif file_name.lower().find('_male_')>=0:
        gender = 'male'
    else:
        print ('gender not found')
        gender = -1
    return gender

def find_last_name(file_name):
    """ get the last name by the rule '__XXXX_birthday_(Female|Male)'
    """
    try:
        found = re.search("(__).+[_](\d+)[_](Female|Male)", file_name)
        last_name_str = found.group(0)
    except AttributeError:
        print('last name number not found')
        return -1
    if find_gender(file_name) == 'male':
        last_name = last_name_str[2:-14]
    elif find_gender(file_name) == 'female':
        last_name = last_name_str[2:-16]
    else:
        print('last name number not found')
        return -1
    return last_name

def find_H_number(file_name):
    """ get the H number by the rule '-Hnumber-'
    """
    
    last_name = find_last_name(file_name)
    try:
        found = re.search("[_|-][H](\d+)", last_name)
        H_str = found.group(0)
    except AttributeError:
        print('H number not found')
        return -1
    H_number = H_str[1:]
    return H_number

def find_birthday(file_name):
    """ get the birthday by '_number_Female' or '_number_Male'
    """
    # still have issues
    try:
        found = re.search("[_](\d+)[_](Female|Male)", file_name)
        birthday_str = found.group(0)
    except AttributeError:
        print('birthday not found')
        return -1
    birthday = birthday_str[1:9]
    return birthday

def find_FOV(file_name):
    """get the field of veiw by "number mm_"
    """
    # find the HD 16mm mode
    if file_name.find('16mm'):
        FOV = 16
        return FOV
    
    try:
        found = re.search("[x](\d+)[m][m]", file_name)
        FOV_str = found.group(0)
    except AttributeError:
        print('field of veiw not found')
        return -1
    FOV = int(FOV_str[1:-2])
    return FOV

def find_OD_OS(file_name):
    """ get eye position by '_OD_' or '_OS_'
    """
    if file_name.upper().find('_OD_')>=0:
        OD_OS = 'OD'
    elif file_name.upper().find('_OS_')>=0:
        OD_OS = 'OS'
    else:
        print ('eye position not found')
        OD_OS = -1
    return OD_OS

def find_scan_time(file_name):
    """ get scanning time by '_number_OD_' or '_number_OS_'
    """
    try:
        found = re.search("[_](\d+)[_](OS|OD)[_]", file_name)
        scantime_str = found.group(0)
    except AttributeError:
        print('scaning time not found')
        return -1
    scantime = scantime_str[1:-4]
    return scantime

def find_save_time(file_name):
    """ get saving time by '_OD_number_' or '_OS_number_'
    """
    try:
        found = re.search("[_](OS|OD)[_](\d+)[_]", file_name)
        savetime_str = found.group(0)
    except AttributeError:
        print('save time not found')
        return -1
    savetime = savetime_str[4:-1]
    return savetime

def find_image_modality(file_name):
    """ get image modality by '_OD_number_Angiography_' or '_OS_number_XX_'
    """
    try:
        found = re.search("[_](OS|OD)[_](\d+)[_](Angiography_|Structure_|B-Scan)", file_name)
        image_modality_str = found.group(0)
    except AttributeError:
        print(file_name)
        print('image modality not found')
        return -1   
    # get the image modality from the re result
    if image_modality_str.lower().find('b-scan')>=0:
        image_modality = 'B-scan'
    elif image_modality_str.lower().find('angiography')>=0:
        image_modality = 'Angiography'
    elif image_modality_str.lower().find('structure')>=0:
        image_modality = 'Structure'
    return image_modality

def find_image_layer(file_name):
    """ get image layer by '_OD|OS_number_Angiography_XX.' 
    """
    try:
        found = re.search("[_](OS|OD)[_](\d+)[_](Angiography_|Structure_|B-Scan).+\.", file_name)
        image_layer_str = found.group(0)
    except AttributeError:
        print('image layer not found')
        return -1
    if find_image_modality(file_name) == 'B-scan':
        image_layer = image_layer_str[19:-1]
    elif find_image_modality(file_name) == 'Angiography':
        image_layer = image_layer_str[31:-1]
    elif find_image_modality(file_name) == 'Structure':
        image_layer = image_layer_str[29:-1]
    else:
        print('image layer not found')
        return -1
    return image_layer
    
def find_file_type(file_name):
    """ get file type by '.XXX'
    """
    try:
        found = re.search("(\.png|\.jpg|\.jpeg|\.tiff|\.tif|\.bmp)", file_name)
        file_type_str = found.group(0)
    except AttributeError:
        print('file type not found')
        return -1
    file_type = file_type_str[1::]
    return file_type
    
    


#%%

def analysis_filename(file_name, ind):
    """splict the filename to a dataframe 
    """
    file_path = file_name
    file_name = clean_file_name(file_name)
    import pandas as pd
    # create the dataframe
    d = {'system_type': find_system_type(file_name),
         'first_name': find_first_name(file_name),
         'case_number': find_case_number(file_name),         
         'last_name': find_last_name(file_name),
         'H_number': find_H_number(file_name),
         'gender': find_gender(file_name),
         'birthday': find_birthday(file_name),
         'FOV': find_FOV(file_name),
         'OD_OS': find_OD_OS(file_name),
         'scan_time': find_scan_time(file_name),
         'save_time': find_save_time(file_name),
         'image_modality': find_image_modality(file_name),
         'image_layer': find_image_layer(file_name),
         'file_type': find_file_type(file_name),
         'file_name': file_name,
         'file_path': file_path
         }
    df = pd.DataFrame(data=d, index=[ind, ])

    return df
    

def analysis_filelists(file_lists):
    """this function accept the file name and 
    """
    
    # make a dataframe to contain the arributes of files 
    file_num = 0;
    for file_name in file_lists:
#        file_name = clean_file_name(file_name)
            
        if file_num==0:
            df = analysis_filename(file_name, file_num)
        else: 
            df = df.append(analysis_filename(file_name, file_num))
        file_num += 1
        
    # save the list sheet
    df.to_csv('file_list.csv')
    return df

#        find_system_type(file_name)
#        find_first_name(file_name)
#        find_case_number(file_name)
#        find_gender(file_name)
#        find_last_name(file_name)
#        find_H_number(file_name)
#        find_birthday(file_name)
#        find_FOV(file_name)
#        find_OD_OS(file_name)
#        find_scan_time(file_name)
#        find_save_time(file_name)
#        find_image_modality(file_name)
#        find_image_layer(file_name)
#        find_file_type(file_name)

def select_by_feature(df, feature_type, value):
    """return the images from the list by a feature
    """
    df_select = df.loc[df[feature_type] == value]
    return df_select

def filter_images_by_features(df, system_type=0, OD_OS=0, FOV=0, 
                             scan_time=0, image_modality=0, image_layer=0 ):
    """Find the images by input features
    Args：
    
    Return:
        the dataframe of one or more images
    """
    # for find one image, the effcient could be very low
    df_new = df
    if system_type!=0:
        df_new = select_by_feature(df_new, 'system_type', system_type)
    if OD_OS!=0:
        df_new = select_by_feature(df_new, 'OD_OS', OD_OS)
    if FOV!=0:
        df_new = select_by_feature(df_new, 'FOV', FOV)
    if scan_time!=0:
        df_new = select_by_feature(df_new, 'scan_time', scan_time)
    if image_modality!=0:
        df_new = select_by_feature(df_new, 'image_modality', image_modality)
    if image_layer!=0:
        df_new = select_by_feature(df_new, 'image_layer', image_layer)       
        
    return df_new
    
   
def get_images_information(df):
    """
    get the unique informations from the dataframe
    """
    unique_scan_time = df['scan_time'].unique()
    unique_FOV = df['FOV'].unique()
    unique_OD_OS = df['OD_OS'].unique()
    unique_case_number = df['case_number'].unique()
    unique_H_number = df['H_number'].unique()
    unique_system_type = df['system_type'].unique()
    unique_gender = df['gender'].unique()
    
    images_information = {'scan_time':unique_scan_time,
                          'FOV': unique_FOV,
                          'OD_OS': unique_OD_OS,
                          'case_number': unique_case_number,
                          'H_number': unique_H_number,
                          'system_type': unique_system_type,
                          'gender': unique_gender
                          }
    return images_information

def get_the_title(df, scan_time):
    """
    get the title of the slide
    """
    df_new = select_by_feature(df, 'scan_time', scan_time)
    # get the informaiton of this scan
    try:
        OD_OS = df_new.iloc[0]['OD_OS']
        FOV = df_new.iloc[0]['FOV']
    except:
        print('cannot get name')
        
    title = scan_time + '_'+ OD_OS + '_' + str(FOV) +'mm'
    return title

def build_the_slide(prs, df, style_type, scan_time):
    img_paths = []
    img_names = []
    if style_type=='Retina':
        #insert the retina slide here
        img_nums = 8

        image_col = ['VRI', 'Superficial', 'Deep', 'Retina']
        image_row = ['Angiography','Structure']
        for row in range(0, 2):
            for col in range(0, 4):
                img_names.append(image_col[col])
                df_img = filter_images_by_features(df, scan_time=scan_time, 
                                          image_modality=image_row[row], 
                                          image_layer=image_col[col])
#                print(df_img)
                path = df_img['file_path'].tolist()
#                print(len(path))
                #if len(path) > 1 print('find mutiple images')
                try:
                    img_paths.append(path[0])
                except IndexError:
                    print('image ' + image_col[col] + 'not found')
                    img_paths.append(-1)
        
        if len(img_paths)<=img_nums:
            title = get_the_title(df, scan_time)
#            print(img_paths)
            prs = add_new_slide_pics(prs, img_paths, img_nums, img_names, title)
        else:
            print ('img_paths have ', len(img_paths), 'larger than 8')
        
        
    elif  style_type=='CC':
        #insert another slide here
        img_nums = 8

        image_col = ['Avascular', 'Custom', 'Choroid', 'RetinaDepthEncoded']
        image_row = ['Angiography','Structure']
        for row in range(0, 2):
            for col in range(0, 4):
                # the last image is B-scanflow
                if ((row==1) & (col==3)):
                    image_col[col] = 'B-ScanFlow'
                    image_row[row] = 'B-scan'
#                    print image_col, image_row
                
                img_names.append(image_col[col])
                df_img = filter_images_by_features(df, scan_time=scan_time, 
                                          image_modality=image_row[row], 
                                          image_layer=image_col[col])
                
                path = df_img['file_path'].tolist()
#                print(path)
                try:
                    img_paths.append(path[0])
                except IndexError:
                    print('image ' + image_col[col] + 'not found')
                    img_paths.append(-1)
                    
                
        if len(img_paths)<=img_nums:
            title = get_the_title(df, scan_time)
            prs = add_new_slide_pics(prs, img_paths, img_nums, img_names, title)
        else:
            print ('img_paths have ', len(img_paths), 'larger than 8')
            
    return prs

def create_slides(df, prs):
    """
    Creat slides
    """
    images_information = get_images_information(df)
    
    scan_times = images_information['scan_time']
    ind_scan = 1
    for scan_time in scan_times:     
        prs = build_the_slide(prs, df, 'Retina', scan_time)
        print ('build the slide for case'+str(ind_scan) + '  for Retina layer')
        prs = build_the_slide(prs, df, 'CC', scan_time)
        print ('build the slide for case'+str(ind_scan) + ' for CC layer')
        
        print ('completed ' + str(ind_scan) + ' of ' + str(len(scan_times)))
        ind_scan += 1
    
    print ('saving...')
    return prs
#%%
######################################################
#%%    
# add table


def add_a_table(prs, rows, cols, slide_title, df):
    """ Add the slide with table
    """
    # 
#    from pptx.util import Inches
# rows = 11, cols = 2
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    
    shapes.title.text = slide_title
    
#    rows = cols = 2
    top = Inches(1.5)
    left = Inches(2.0)
    width = Inches(6.0)
    height = Inches(0.8)
    
    table = shapes.add_table(rows, cols, left, top, width, height).table
    
#    # set column widths
#    table.columns[0].width = Inches(2.0)
#    table.columns[1].width = Inches(4.0)

    # get images information
    file_info = get_images_information(df)
#     write column headings
    table.cell(0, 0).text = 'Patient ID'
    table.cell(0, 1).text = str(file_info['case_number'])
    
    table.cell(1, 0).text = 'MRN'
    table.cell(1, 1).text = str(file_info['H_number'])
    
    table.cell(2, 0).text = 'Gender'
    table.cell(2, 1).text = str(file_info['gender'])
    
    table.cell(3, 0).text = 'Age'
    table.cell(4, 0).text = 'Pathology'
    table.cell(5, 0).text = 'Clinical image'
    
    table.cell(6, 0).text = 'Image eyes'
    table.cell(6, 1).text = str(file_info['OD_OS'])
    
    table.cell(7, 0).text = 'Image date'
    table.cell(7, 1).text = str(file_info['scan_time'])
    
    table.cell(8, 0).text = 'Treatment times'
    table.cell(9, 0).text = 'Image times'
    
    table.cell(10, 0).text = 'Image system'
    table.cell(10, 1).text = str(file_info['system_type'])
    
    table.cell(11, 0).text = 'Scanning protocol'
    table.cell(11, 1).text = str(file_info['FOV'])
    
    
    
#    
#    # write body cells
#    table.cell(1, 0).text = 'Baz'
#    table.cell(1, 1).text = 'Qux'
    return prs