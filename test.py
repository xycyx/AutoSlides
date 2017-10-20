#!/usr/bin/env python2
# -*- coding: utf-8 -*-
"""
Created on Thu Oct 12 16:14:01 2017

@author: yxcheng
"""


#%%
from pptx import Presentation
prs = Presentation()
#blank_slide_layout = prs.slide_layouts[6]
#slide = prs.slides.add_slide(blank_slide_layout)
#img_path = './example/Picture4.png'
#add_pic_with_title(slide, img_path, 5, 4.5, 2.5, 'Picture4.png')
#prs.save('test2.pptx')

import os
img_paths = []
for file in os.listdir("./example"):
    if file.endswith(".png"):
        print(os.path.join("./example", file))
        img_paths.append(os.path.join("./example", file))



img_name = ['0','1','2','3','4','5','6','7']
img_nums = 6
prs = add_new_slide_pics(prs, img_paths, img_nums, img_name)

add_a_table(prs, 4, 3, 'test the table')
prs.save('test1.pptx')
#%%
def add_pic_with_title(slide, img_path, left, top, width, img_name):
    """add a picture to slide from img_path and add the title with img_name
    """
    from pptx.util import Inches
    left = Inches(left)
    top_img = Inches(top)
    width = Inches(width)
    
    # add picture
    pic = slide.shapes.add_picture(img_path, left, top_img, width=width)
    
    #add title
    top_title = Inches(top - 0.3)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top_title , width, height)
    tf = txBox.text_frame
    tf.text = img_name
    
    return slide
    
def add_new_slide_pics(prs, img_paths, img_nums, img_name):
    """Decide the layout of pictures
    """
#    from pptx.util import Inches
    #img_nums = len(img_paths)

    #%% 2
    if img_nums == 2:
        layout = [1, 2]
        prs = set_layout(prs, layout, img_paths, img_nums, img_name)
    
    #%% 3
    elif img_nums == 3:
        layout = [1, 3]
        prs = set_layout(prs, layout, img_paths, img_nums, img_name)
   
    #%% 4 still have a issue
    elif img_nums == 4:
        layout = [2, 2]
        prs = set_layout(prs, layout, img_paths, img_nums, img_name)
    
    #%% 6
    elif img_nums == 6:
        layout = [2, 3]
        prs = set_layout(prs, layout, img_paths, img_nums, img_name)
    
    #%% 8
    elif img_nums == 8:
        layout = [2, 4]
        prs = set_layout(prs, layout, img_paths, img_nums, img_name)        
               
    else:
        print('other numbers')
    return prs
    
def set_layout(prs, layout,img_paths, img_nums, img_name):
    """Set the position and decide the layout of picture 
    Args: 
        layout: the size of images matriax , as [4, 2]
    Return:
        pre: the Presentation 
    """
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    ind = 0
    int_width = 0.3
    width = (10 - int_width)/layout[1]
    pic_width = width - int_width
    int_top = 7.5 - layout[0]*width
#    for img_path in img_paths: 
    for i in range(0, img_nums):
        img_path = img_paths[i]          
        col = ind%layout[1]
        row = ind//layout[1]
#            height = ((10 - int_width)/layout[1] - int_width
        left = col*width + int_width
        top = row*width + int_top
        
        # add the picture on its position
        add_pic_with_title(slide, img_path, left, top, pic_width, img_name[ind])
        ind = ind + 1
    return prs
 #%%   
 # get the size of image without open it
from PIL import Image
im = Image.open('./example/1.tiff')
im.size


#width, height = get_image_size('./example/1.jpg')
#import struct
#import imghdr
#
#def get_image_size(fname):
#    '''Determine the image type of fhandle and return its size.
#    from draco'''
#    with open(fname, 'rb') as fhandle:
#        head = fhandle.read(24)
#        if len(head) != 24:
#            return
#        if imghdr.what(fname) == 'png':
#            check = struct.unpack('>i', head[4:8])[0]
#            if check != 0x0d0a1a0a:
#                return
#            width, height = struct.unpack('>ii', head[16:24])
#        elif imghdr.what(fname) == 'gif':
#            width, height = struct.unpack('<HH', head[6:10])
#        elif imghdr.what(fname) == 'bmp':
#            _, width, height, depth = re.search(
#                b"((\d+)\sx\s"
#                b"(\d+)\sx\s"
#                b"(\d+))", str).groups()
#            width = int(width)
#            height = int(height)
#        elif imghdr.what(fname) == 'jpeg':
#            try:
#                fhandle.seek(0) # Read 0xff next
#                size = 2
#                ftype = 0
#                while not 0xc0 <= ftype <= 0xcf:
#                    fhandle.seek(size, 1)
#                    byte = fhandle.read(1)
#                    while ord(byte) == 0xff:
#                        byte = fhandle.read(1)
#                    ftype = ord(byte)
#                    size = struct.unpack('>H', fhandle.read(2))[0] - 2
#                # We are at a SOFn block
#                fhandle.seek(1, 1)  # Skip `precision' byte.
#                height, width = struct.unpack('>HH', fhandle.read(4))
#            except Exception: #IGNORE:W0703
#                return
#        else:
#            return
#        return width, height    
#%%    
# add table


def add_a_table(prs, rows, cols, slide_title):
    """ Add the slide with table
    """
    from pptx.util import Inches
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    
    shapes.title.text = slide_title
    
#    rows = cols = 2
    left = top = Inches(2.0)
    width = Inches(6.0)
    height = Inches(0.8)
    
    table = shapes.add_table(rows, cols, left, top, width, height).table
    
#    # set column widths
#    table.columns[0].width = Inches(2.0)
#    table.columns[1].width = Inches(4.0)
#    
#    # write column headings
#    table.cell(0, 0).text = 'Foo'
#    table.cell(0, 1).text = 'Bar'
#    
#    # write body cells
#    table.cell(1, 0).text = 'Baz'
#    table.cell(1, 1).text = 'Qux'
    return prs
#######################################################################

def search_images_in_folder(directory):
    """Search and list all images (bmp, png, jpg, tif) files in the folder
    Args: 
        directory: the folder conatins all files
    Return:
        file_lists: the list of all images files
    """
    file_lists = []
#    directory = './example2'
    num = 0
    for root, dirs, files in os.walk(directory):
        for file in files:
            if file.lower().endswith(('.bmp','.png','.jpg','.jpeg','.tiff','.tif')):
#                print file
                file_lists.append(os.path.join(root, file))
    return file_lists       
 
file_lists = search_images_in_folder('./example2')

#%%            
import re
file_name = file_lists[0]
file_name = clean_file_name(file_name)

file_name2 = 'UWEVEREST_981__UWEVEREST-981-H2212406_19750320_Male_Angio (3mmx3mm)_20170104135911_OS_20170104144338_Angiography_Avascular.bmp'
file_name2 = clean_file_name(file_name2)

file_name3 = './example2/Angio6x6/UWANG5000_977__UWANG5000-977-H3848027_19741016_Female_Angiography 6x6 mm_20161123162131_OS_20161123163020_Structure_Superficial.bmp'
file_name3 = clean_file_name(file_name3)

try:
    found = re.search("(?<=Angiography).*?(?=mm)", file_name)
except AttributeError:
    print(' pattern not found')
#    return -1
found.group(0) 


def clean_file_name(file_name):
    """Remove the ' ', '(', ')' of the the filename
    Args: 
        file_name: the input filename
    Returns: 
        file_name: the output of cleaned filename
    """
    
    while ' ' in file_name:
        file_name = file_name.replace(' ', '')
    while '(' in file_name:
        file_name = file_name.replace('(', '')
    while ')' in file_name:
        file_name = file_name.replace(')', '')           
    return file_name

#%%
# regular expression to parse the file name
def find_system_type(file_name):
    """ get system type by the rule 'UWANG5000' or 'UWEVEREST'
    """
    if file_name.upper().find('UWANG5000')>=0:
        system_type = 'SD'
    elif file_name.upper().find('UWEVEREST')>=0:
        system_type = 'SS'
    else:
        print ('system type not found')
        system_type = -1
    return system_type

def find_first_name(file_name):
    """ get the first name by the rule '_XXXX__'
    """
    # the the last occurrence of '/', as the strating position of firstname
    if file_name.rfind('/')>=0:
        name_position = file_name.rfind('/') + 1
    else:
        name_position = 0
    
    #get the first name from the strat position    
    try:
        found = re.search(".+[_][_]", file_name[name_position:])
        first_name_str = found.group(0)
    except AttributeError:
        print('first name number not found')
        return -1
    first_name = first_name_str[:-2]
    return first_name

def find_case_number(file_name):
    """ get the case number by rule '-number-'
    """
    case_str = find_first_name(file_name)
    try:
        found = re.search("[-|_]\d+", case_str)
        case_str_number = found.group(0)
    except AttributeError:
        print('case number not found')
        return -1
    case_number = int(case_str_number[1:])
    return case_number

def find_gender(file_name):
    """ get system type by '_male_' or '_female_'
    """
    if file_name.lower().find('_female_')>=0:
        gender = 'female'
    elif file_name.lower().find('_male_')>=0:
        gender = 'male'
    else:
        print ('gender not found')
        gender = -1
    return gender

def find_last_name(file_name):
    """ get the last name by the rule '__XXXX_birthday_(Female|Male)'
    """
    try:
        found = re.search("(__).+[_](\d+)[_](Female|Male)", file_name)
        last_name_str = found.group(0)
    except AttributeError:
        print('last name number not found')
        return -1
    if find_gender(file_name) == 'male':
        last_name = last_name_str[2:-14]
    elif find_gender(file_name) == 'female':
        last_name = last_name_str[2:-16]
    else:
        print('last name number not found')
        return -1
    return last_name

def find_H_number(file_name):
    """ get the H number by the rule '-Hnumber-'
    """
    
    last_name = find_last_name(file_name)
    try:
        found = re.search("[_|-][H](\d+)", last_name)
        H_str = found.group(0)
    except AttributeError:
        print('H number not found')
        return -1
    H_number = H_str[1:]
    return H_number

def find_birthday(file_name):
    """ get the birthday by '_number_Female' or '_number_Male'
    """
    # still have issues
    try:
        found = re.search("[_](\d+)[_](Female|Male)", file_name)
        birthday_str = found.group(0)
    except AttributeError:
        print('birthday not found')
        return -1
    birthday = birthday_str[1:9]
    return birthday

def find_FOV(file_name):
    """get the field of veiw by "number mm_"
    """
    try:
        found = re.search("(\d+)[m][m][_]", file_name)
        FOV_str = found.group(0)
    except AttributeError:
        print('field of veiw not found')
        return -1
    FOV = int(FOV_str[0:-3])
    return FOV

def find_OD_OS(file_name):
    """ get eye position by '_OD_' or '_OS_'
    """
    if file_name.upper().find('_OD_')>=0:
        OD_OS = 'OD'
    elif file_name.upper().find('_OS_')>=0:
        OD_OS = 'OS'
    else:
        print ('eye position not found')
        OD_OS = -1
    return OD_OS

def find_scan_time(file_name):
    """ get scanning time by '_number_OD_' or '_number_OS_'
    """
    try:
        found = re.search("[_](\d+)[_](OS|OD)[_]", file_name)
        scantime_str = found.group(0)
    except AttributeError:
        print('scaning time not found')
        return -1
    scantime = scantime_str[1:-4]
    return scantime

def find_save_time(file_name):
    """ get saving time by '_OD_number_' or '_OS_number_'
    """
    try:
        found = re.search("[_](OS|OD)[_](\d+)[_]", file_name)
        savetime_str = found.group(0)
    except AttributeError:
        print('save time not found')
        return -1
    savetime = savetime_str[4:-1]
    return savetime

def find_image_modality(file_name):
    """ get image modality by '_OD_number_Angiography_' or '_OS_number_XX_'
    """
    try:
        found = re.search("[_](OS|OD)[_](\d+)[_](Angiography_|Structure_)", file_name)
        image_modality_str = found.group(0)
    except AttributeError:
        print('image modality not found')
        return -1   
    # get the image modality from the re result
    if image_modality_str.lower().find('angiography')>=0:
        image_modality = 'Angiography'
    elif image_modality_str.lower().find('structure')>=0:
        image_modality = 'Structure'
    return image_modality

def find_image_layer(file_name):
    """ get image layer by '_OD|OS_number_Angiography_XX.' 
    """
    try:
        found = re.search("[_](OS|OD)[_](\d+)[_](Angiography_|Structure_).+\.", file_name)
        image_layer_str = found.group(0)
    except AttributeError:
        print('image layer not found')
        return -1
    if find_image_modality(file_name) == 'Angiography':
        image_layer = image_layer_str[31:-1]
    elif find_image_modality(file_name) == 'Structure':
        image_layer = image_layer_str[29:-1]
    else:
        print('image layer not found')
        return -1
    return image_layer
    
def find_file_type(file_name):
    """ get file type by '.XXX'
    """
    try:
        found = re.search("(\.png|\.jpg|\.jpeg|\.tiff|\.tif|\.bmp)", file_name)
        file_type_str = found.group(0)
    except AttributeError:
        print('file type not found')
        return -1
    file_type = file_type_str[1::]
    return file_type
    
    


#%%

def analysis_filename(file_name):
    """splict the filename to a dataframe 
    """
    import pandas as pd
    d = {'system_type': find_system_type(file_name),
         'first_name': find_first_name(file_name),
         'case_number': find_case_number(file_name),
         
         'last_name': find_last_name(file_name),
         'H_number': find_H_number(file_name),
         'gender': find_gender(file_name),
         'birthday': find_birthday(file_name),
         'FOV': find_FOV(file_name),
         'OD_OS': find_OD_OS(file_name),
         'scan_time': find_scan_time(file_name),
         'save_time': find_save_time(file_name),
         'image_modality': find_image_modality(file_name),
         'image_layer': find_image_layer(file_name),
         'file_type': find_file_type(file_name),
         'file_name': file_name
         }
    df = pd. DataFrame(data=d, index=[0,])

    return df
    

def analysis_filelists(file_lists):
    """this function accept the file name and 
    """
    
    # make a dataframe to contain the arributes of files 
    df = pd.
    for file_name in file_lists:
        file_name = clean_file_name(file_name)
        analysis_filename(file_name)    
        
#        find_system_type(file_name)
#        find_first_name(file_name)
#        find_case_number(file_name)
#        find_gender(file_name)
#        find_last_name(file_name)
#        find_H_number(file_name)
#        find_birthday(file_name)
#        find_FOV(file_name)
#        find_OD_OS(file_name)
#        find_scan_time(file_name)
#        find_save_time(file_name)
#        find_image_modality(file_name)
#        find_image_layer(file_name)
#        find_file_type(file_name)
        
                
######################################################
#%% test part
from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!" 
subtitle.text = "python-pptx was here!"

prs.save('test.pptx')

#%%
from pptx import Presentation

prs = Presentation()
bullet_slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Adding a Bullet Slide'

tf = body_shape.text_frame
tf.text = 'Find the bullet slide layout'

p = tf.add_paragraph()
p.text = 'Use _TextFrame.text for first bullet'
p.level = 1

p = tf.add_paragraph()
p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
p. level = 2

prs.save('test.pptx')
#%%
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = width = height = Inches(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame

tf.text = "This is text inside a textbox"

p = tf.add_paragraph()
p.text = "This is a second paragraph that's bold"
p.font.bold = True

p = tf.add_paragraph()
p.text = "This is a third paragraph that's big"
p.font.size = Pt(40)

prs.save('test.pptx')
#%%
from pptx import Presentation
from pptx.util import Inches



prs = Presentation() 
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

img_path = './example/Picture1.png'

left = top = Inches(1)
pic = slide.shapes.add_picture(img_path, left, top)

left = top = width = height = Inches(1)
txBox = slide.shapes.add_textbox(left, top*0.7 , width, height)
tf = txBox.text_frame

tf.text = img_path[10::]


img_path = './example/Picture2.png'
left = Inches(5)
height = Inches(2.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

top = width = height = Inches(1)
txBox = slide.shapes.add_textbox(left, top*0.7 , width, height)
tf = txBox.text_frame

tf.text = img_path[10::]
#%%
img_path = './example/Picture2.png'
left = Inches(5)
height = Inches(2.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

top = width = height = Inches(1)
txBox = slide.shapes.add_textbox(left, top*0.7 , width, height)
tf = txBox.text_frame

tf.text = img_path[10::]
    
#%% main
import argparse
if __name__ == "__main__":
    args = parse_args()
    
    